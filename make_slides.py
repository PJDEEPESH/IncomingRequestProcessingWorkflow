"""
Generates the 5-slide submission deck (POC_Slides.pptx) in the fixed structure
the brief requires. Clean, light, corporate — themed to Firstsource's identity
(navy + white + blue, with a warm accent).

Run:  .venv\\Scripts\\python.exe make_slides.py
Then: edit the DEMO_LINK below, and drop Firstsource's official logo into the
top-left placeholder (Insert > Picture) if you have it.
"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

DEMO_LINK = "https://liberation-production.up.railway.app"
REPO_LINK = "<< optional: GitHub repo link >>"

# Firstsource-inspired corporate palette (light theme)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NAVY = RGBColor(0x0A, 0x1A, 0x33)      # headings / text
BLUE = RGBColor(0x0B, 0x5F, 0xD6)      # primary accent
CORAL = RGBColor(0xFF, 0x5A, 0x3C)     # highlight accent
GRAY = RGBColor(0x5B, 0x64, 0x72)      # muted text
PANEL = RGBColor(0xF2, 0xF5, 0xFA)     # box fill
LINE = RGBColor(0xD9, 0xE0, 0xEC)
GREEN = RGBColor(0x1F, 0x9D, 0x57)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def new_slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    # brand lockup (top-left) — replace the square with the official logo if you have it
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(0.42), Inches(0.22), Inches(0.22))
    dot.fill.solid(); dot.fill.fore_color.rgb = CORAL; dot.line.fill.background()
    _text(s, "Firstsource", 0.98, 0.38, 3, 0.35, size=14, color=NAVY, bold=True)
    # footer
    _text(s, "Incoming Request Processing Workflow  ·  POC  ·  Deepesh",
          0.7, 7.05, 9, 0.3, size=9, color=GRAY)
    return s


def _text(slide, txt, left, top, width, height, size=18, color=NAVY, bold=False,
          align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Segoe UI"
    return tb


def bullets(slide, items, left, top, width, height, size=15, gap=9):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        txt, col = item if isinstance(item, tuple) else (item, NAVY)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run()
        r.text = "▪   " + txt
        r.font.size = Pt(size)
        r.font.color.rgb = col
        r.font.name = "Segoe UI"
    return tb


def header(slide, kicker, title):
    _text(slide, kicker, 0.7, 1.15, 11, 0.35, size=13, color=BLUE, bold=True)
    _text(slide, title, 0.7, 1.5, 12, 0.8, size=28, color=NAVY, bold=True)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(2.22), Inches(0.9), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()


def box(slide, txt, left, top, w, h, fill=PANEL, tcolor=NAVY, size=11):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = LINE; sh.line.width = Pt(1)
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.size = Pt(size); r.font.color.rgb = tcolor; r.font.bold = True; r.font.name = "Segoe UI"
    return sh


def glyph(slide, ch, left, top, w=0.5, h=0.35, size=16):
    _text(slide, ch, left, top, w, h, size=size, color=GRAY, align=PP_ALIGN.CENTER)


# ---------------- Slide 1 ----------------
s = new_slide()
header(s, "SLIDE 1  ·  01", "Problem Understanding and Objective")
_text(s, "The problem", 0.7, 2.5, 6, 0.4, size=16, color=BLUE, bold=True)
bullets(s, [
    "Support / BPO teams receive high volumes of mixed, unlabelled messages across email, web forms, chat and shared inboxes.",
    "A human must read each one, judge type and urgency, route it, and log it — this manual triage is slow, inconsistent, and costly to scale.",
    "Urgent complaints wait in the same queue as trivial questions, so SLAs slip and cost per contact stays high.",
], 0.7, 2.95, 5.9, 3)
_text(s, "The objective", 7.0, 2.5, 6, 0.4, size=16, color=BLUE, bold=True)
bullets(s, [
    "An AI prototype that classifies each request AND runs a distinct multi-step remediation workflow per type — not just a label.",
    "Prioritise urgent cases, draft the reply, route to the right team, and hand off to a human when unsure.",
], 7.0, 2.95, 5.6, 3)
_text(s, "Directly maps to Firstsource's core BPM business — customer operations across Communications, Banking & Healthcare.",
      0.7, 6.4, 12, 0.5, size=12, color=CORAL, bold=True)

# ---------------- Slide 2 ----------------
s = new_slide()
header(s, "SLIDE 2  ·  02", "Solution Architecture and Design Flow")
row1 = 2.6
box(s, "Incoming Request\nEmail · Web Form · Chat", 0.75, row1, 2.5, 0.9)
glyph(s, "→", 3.3, row1 + 0.28)
box(s, "AI Classify\ntype · urgency · sentiment", 3.85, row1, 2.5, 0.9, fill=RGBColor(0xE2, 0xEE, 0xFC))
glyph(s, "→", 6.4, row1 + 0.28)
box(s, "AI Reason\n(thinking node)", 6.95, row1, 2.5, 0.9, fill=RGBColor(0xEC, 0xE7, 0xFB))
glyph(s, "→", 9.5, row1 + 0.28)
box(s, "Router\nconfidence gate", 10.05, row1, 2.5, 0.9, fill=RGBColor(0xFF, 0xF0, 0xE6))
glyph(s, "↓", 6.5, 3.65, size=18)
row2 = 4.05
chips = [("COMPLAINT", RGBColor(0xFF, 0xE8, 0xE0)), ("ENQUIRY", RGBColor(0xE4, 0xF5, 0xEA)),
         ("SERVICE", RGBColor(0xE2, 0xEE, 0xFC)), ("ESCALATION", RGBColor(0xFF, 0xD9, 0xD2)),
         ("HUMAN REVIEW", RGBColor(0xEC, 0xE7, 0xFB))]
cw, gap = 2.2, 0.22
start = (13.333 - (5 * cw + 4 * gap)) / 2
for i, (name, col) in enumerate(chips):
    box(s, name, start + i * (cw + gap), row2, cw, 0.7, fill=col)
glyph(s, "↓", 6.5, 4.85, size=18)
box(s, "Output  +  Case Log (SQLite)  +  Human Approval  +  reply written to /outbox",
    2.67, 5.3, 8.0, 0.7, fill=PANEL, size=12)
_text(s, "System design:   React + TypeScript (UI)   →   FastAPI (API)   →   LangGraph agentic pipeline   →   OpenAI + SQLite",
      0.7, 6.45, 12, 0.5, size=12.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

# ---------------- Slide 3 ----------------
s = new_slide()
header(s, "SLIDE 3  ·  03", "Implementation Highlights")
bullets(s, [
    ("Agentic LangGraph state machine with TWO AI reasoning nodes — a Classify node and a Reason (thinking) node — then a conditional router into 5 branches.", NAVY),
    ("Confidence-gated human-in-the-loop — if the AI is < 60% sure, the case is held for human approval instead of being guessed.", NAVY),
    ("Sentiment detection + AI-chosen department routing (Billing / Provisioning / Field / Network Ops).", NAVY),
    ("Graceful degradation — a rule-based fallback keeps the system running, with a visible warning, if the AI API is unavailable.", NAVY),
    ("Measured quality — 87% type-classification accuracy on a labelled test set, shown in-app with a confusion matrix.", GREEN),
], 0.7, 2.5, 12, 3.4)
_text(s, "Code: LangGraph nodes & edges in workflows.py  ·  prompts in classifier.py and responder.py",
      0.7, 6.4, 12, 0.4, size=12, color=GRAY)

# ---------------- Slide 4 ----------------
s = new_slide()
header(s, "SLIDE 4  ·  04", "Challenges and Learnings")
bullets(s, [
    ("Challenge: telling an angry Complaint from a true Escalation — solved with tone-aware prompting and a confidence gate that routes uncertain cases to a human.", NAVY),
    ("Challenge: corporate SSL inspection blocked the OpenAI API — solved by trusting the OS certificate store (truststore), without disabling security.", NAVY),
    ("Trade-off: LangGraph vs plain functions — chose LangGraph for explicit, auditable branching and a clean architecture diagram.", NAVY),
    ("Honest limitation: the inbox is simulated and some downstream steps are stubs that map to Salesforce / Twilio / email in production; the reply action is really executed.", GRAY),
    ("Learning: reliability — a fallback engine, output validation, and a measured accuracy number — matters as much as the AI itself.", GREEN),
], 0.7, 2.5, 12, 3.5)

# ---------------- Slide 5 ----------------
s = new_slide()
header(s, "SLIDE 5  ·  05", "Demo Summary and Next Steps")
_text(s, "Live demo:  " + DEMO_LINK, 0.7, 2.5, 12, 0.4, size=14, color=BLUE, bold=True)
_text(s, "Repository:  " + REPO_LINK, 0.7, 2.92, 12, 0.35, size=12, color=GRAY)
_text(s, "What it delivers", 0.7, 3.5, 6, 0.4, size=15, color=BLUE, bold=True)
bullets(s, [
    "Multi-source intake → AI classify + reason → 5 branch workflows.",
    "Human-approval safety net, full audit log, KPI dashboard.",
    "87% measured accuracy; runs even if the AI API is down.",
], 0.7, 3.9, 5.9, 3, size=14)
_text(s, "Next steps", 7.0, 3.5, 6, 0.4, size=15, color=BLUE, bold=True)
bullets(s, [
    "Add new nodes easily (e.g. fraud / abuse detection) — it's a graph.",
    "Real integrations: ticketing / PagerDuty escalation, live email intake.",
    "Knowledge-base grounding for enquiry answers; real SLA timers.",
    "Feedback loop to improve accuracy; multilingual support.",
], 7.0, 3.9, 5.6, 3, size=14)

prs.save("POC_Slides.pptx")
print("Saved POC_Slides.pptx  (", len(prs.slides._sldIdLst), "slides )")
